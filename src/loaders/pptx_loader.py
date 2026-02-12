from pathlib import Path
from typing import Optional
from pptx import Presentation

from src.loaders.base_loader import BaseLoader
from src.utils.document import Document

"""
PPTXLoader - Load Microsoft PowerPoint presentations
Extracts text from .pptx files slide by slide
"""
"""
Load a Microsoft PowerPoint (.pptx) file and extract text
Each slide can be a separate document or all slides can be
combined into a single document
Attributes:
    file_path: Path to the PowerPoint file
    per_slide: Whether to create one document per slide
Example:
    >>> loader = PPTXLoader("presentation.pptx")
    >>> docs = loader.load()
    >>> print(f"Loaded {len(docs)} slides")
"""
class PPTXLoader(BaseLoader):
    
    """
    Initialize the PPTXLoader
    Args:
        file_path: Path to the PowerPoint file
        per_slide: If True, create one document per slide
                    If False, combine all slides into one document
        include_notes: Whether to include speaker notes
        metadata: Optional additional metadata
    """
    def __init__(
        self,
        file_path: str,
        per_slide: bool = True,
        include_notes: bool = True,
        metadata: Optional[dict] = None
    ):
        self.file_path = Path(file_path)
        self.per_slide = per_slide
        self.include_notes = include_notes
        self.extra_metadata = metadata or {}
    
    """
    Load the PowerPoint file and extract text
    Returns:
        List of Documents (one per slide or single combined)
    Raises:
        FileNotFoundError: If the file doesn't exist
    """
    def load(self) -> list[Document]:
        if not self.file_path.exists():
            raise FileNotFoundError(f"File not found: {self.file_path}")
        
        prs = Presentation(self.file_path)
        total_slides = len(prs.slides)
        
        if self.per_slide:
            return self._load_per_slide(prs, total_slides)
        else:
            return self._load_combined(prs, total_slides)
    
    """
    Load the PowerPoint file and extract text
    Returns:
        List of Documents (one per slide or single combined)
    Raises:
        FileNotFoundError: If the file doesn't exist
    """
    def _load_per_slide(self, prs, total_slides: int) -> list[Document]:
        documents = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            content = self._extract_slide_text(slide, slide_num)
            
            if content.strip():
                metadata = {
                    "source": str(self.file_path),
                    "file_name": self.file_path.name,
                    "file_type": "pptx",
                    "slide": slide_num,
                    "total_slides": total_slides,
                    **self.extra_metadata
                }
                
                documents.append(Document(page_content=content, metadata=metadata))
        
        return documents
    
    """
    Load the PowerPoint file and extract text
    Returns:
        List of Documents (one per slide or single combined)
    Raises:
        FileNotFoundError: If the file doesn't exist
    """
    def _load_combined(self, prs, total_slides: int) -> list[Document]:
        all_content = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            content = self._extract_slide_text(slide, slide_num)
            if content.strip():
                all_content.append(content)
        
        combined_content = "\n\n---\n\n".join(all_content)
        
        metadata = {
            "source": str(self.file_path),
            "file_name": self.file_path.name,
            "file_type": "pptx",
            "total_slides": total_slides,
            **self.extra_metadata
        }
        
        return [Document(page_content=combined_content, metadata=metadata)]
    
    """
    Extract all text from a slide
    Args:
        slide: The slide to extract text from
        slide_num: The slide number
    Returns:
        The extracted text
    """
    def _extract_slide_text(self, slide, slide_num: int) -> str:
        parts = [f"## Slide {slide_num}"]
        
        """
        Extract text from shapes
        """
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
            
            """
            Handle tables
            """
            if shape.has_table:
                table_text = self._extract_table(shape.table)
                if table_text:
                    parts.append(table_text)
        
        """
        Extract speaker notes
        """
        if self.include_notes and slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text.strip():
                parts.append(f"\nNotes: {notes_frame.text.strip()}")
        
        return "\n\n".join(parts)
    
    """
    Extract text from a table
    Args:
        table: The table to extract text from
    Returns:
        The extracted text
    """
    def _extract_table(self, table) -> str:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)
